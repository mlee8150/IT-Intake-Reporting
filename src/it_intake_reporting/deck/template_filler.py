"""Fills the real template (templates/Technology_Intake_Weekly_Exec_Summary.pptx)
with a week's computed values.

Shape names below (e.g. "Text 7") come straight from the template's own XML
— they're what whatever tool exported this template happened to call each
box, not anything meaningful on its own. `FIELD_SHAPES` is the one place
that maps a logical field name to its shape, so if the template gets
regenerated (new shape names) only this map needs updating, not the panel
code that produces the values.

Text is written into the *existing* run(s) of each shape rather than via
`text_frame.text = ...`, which would collapse the box to a single unstyled
run and strip the template's formatting (color, size, font) — see the pptx
skill's editing notes.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE

# logical field name -> template shape name, for single-run text boxes.
FIELD_SHAPES: dict[str, str] = {
    "week_ending": "Text 1",
    "data_as_of": "Text 2",
    "headline_open_active": "Text 7",
    "headline_open_active_delta": "Text 8",
    "headline_resolved": "Text 13",
    "headline_resolved_delta": "Text 14",
    "headline_aging_90": "Text 19",
    "headline_aging_90_delta": "Text 20",
    "headline_exec_critical": "Text 25",
    "headline_exec_critical_delta": "Text 26",
    "p1_opening_balance": "Text 39",
    "p1_newly_opened": "Text 40",
    "p1_completed": "Text 42",
    "p1_not_approved": "Text 44",
    "p1_cancelled": "Text 46",
    "p1_on_hold": "Text 48",
    "p1_re_opened": "Text 50",
    "p2_this_week_opened": "Text 62",
    "p2_this_week_closed": "Text 64",
    "p3_hours_CyberArch": "Text 79",
    "p3_hours_TPRM": "Text 82",
    "p3_hours_Legal": "Text 85",
    "p3_hours_IAM": "Text 88",
    "p3_hours_AI": "Text 91",
    "p3_hours_Sub-ARB": "Text 94",
    "p3_hours_total": "Text 97",
    "p4_median_days": "Text 106",
    "p4_delta": "Text 108",
    "p6_bucket_0_30": "Text 127",
    "p6_bucket_31_60": "Text 130",
    "p6_bucket_61_90": "Text 133",
    "p6_bucket_90_plus": "Text 136",
    "p7_stalled_total": "Text 143",
    "p7_top_team_1_name": "Text 147",
    "p7_top_team_1_count": "Text 148",
    "p7_top_team_2_name": "Text 150",
    "p7_top_team_2_count": "Text 151",
    "p7_top_team_3_name": "Text 153",
    "p7_top_team_3_count": "Text 154",
}

# "Text 55" holds the closing balance as two runs: the number, then the delta
# (e.g. run0="102", run1="  ▲ 2") — handled separately from FIELD_SHAPES
# because it can't be treated as one plain string without losing that split.
CLOSING_BALANCE_SHAPE = "Text 55"

CHART_NAMES = ["Chart 0", "Chart 1", "Chart 2", "Chart 3"]


@dataclass(frozen=True)
class ChartUpdate:
    categories: list[str]
    series: dict[str, list[float]]  # series name -> one value per category


@dataclass
class DeckValues:
    text_fields: dict[str, str]  # keys from FIELD_SHAPES
    closing_balance_value: str
    closing_balance_delta: str
    # keyed "Chart 0".."Chart 3" — see CHART_NAMES / docs/DATA_SOURCES.md for
    # which panel each index is.
    charts: dict[str, ChartUpdate] = field(default_factory=dict)


def fill_deck(template_path: Path, output_path: Path, values: DeckValues) -> None:
    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    shapes_by_name = _index_shapes(slide.shapes)

    missing = set(FIELD_SHAPES) - set(values.text_fields)
    if missing:
        raise ValueError(f"DeckValues is missing text fields: {sorted(missing)}")

    for field_name, text in values.text_fields.items():
        shape_name = FIELD_SHAPES[field_name]
        _set_single_run_text(shapes_by_name[shape_name], text)

    _set_two_run_text(
        shapes_by_name[CLOSING_BALANCE_SHAPE],
        values.closing_balance_value,
        values.closing_balance_delta,
    )

    for chart_name, update in values.charts.items():
        shape = shapes_by_name[chart_name]
        if not shape.has_chart:
            raise ValueError(f"Shape {chart_name!r} is not a chart shape")
        chart_data = CategoryChartData()
        chart_data.categories = update.categories
        for series_name, series_values in update.series.items():
            chart_data.add_series(series_name, series_values)
        shape.chart.replace_data(chart_data)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _index_shapes(shapes) -> dict[str, object]:
    index: dict[str, object] = {}
    for shape in shapes:
        index[shape.name] = shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            index.update(_index_shapes(shape.shapes))
    return index


def _set_single_run_text(shape, text: str) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    if not paragraph.runs:
        paragraph.add_run()
    paragraph.runs[0].text = text
    for extra_run in paragraph.runs[1:]:
        extra_run.text = ""


def _set_two_run_text(shape, run0_text: str, run1_text: str) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    if len(paragraph.runs) < 2:
        raise ValueError(
            f"Expected shape {shape.name!r} to have 2 runs (value + delta), "
            f"found {len(paragraph.runs)}. Template may have changed — see "
            "this module's docstring."
        )
    paragraph.runs[0].text = run0_text
    paragraph.runs[1].text = run1_text
